import collections
import collections.abc
import pptx
from pptx.util import Inches
from jira import JIRA
import configparser
import sys

if len(sys.argv) == 4:
    start_quarter = int(sys.argv[1])
    end_quarter = int(sys.argv[2])
    strategic_theme = sys.argv[3]
else:
    start_quarter = 1
    end_quarter = 4
    strategic_theme = "T-Adoption to Retention" 

config = configparser.ConfigParser()
config.read("config.ini")

jira_server = config.get("jira", "server")
jira_username = config.get("jira", "username")
jira_password = config.get("jira", "password")

jira = JIRA(server=jira_server, basic_auth=(jira_username, jira_password))


jql = f'project = SP and "Investment Cluster & Strategic Theme" = "{strategic_theme}" and status != Obsolete and issuetype = "Portfolio Epic"'
params = {"jql": jql}


issues = jira.search_issues(jql, maxResults=5000)


def create_slide(prs):
    slide_layout = prs.slide_layouts[6]  # Use the blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    return slide


def create_table(slide, num_rows, num_cols):
    top_padding = Inches(0.2)  # Set top padding value
    table_top = top_padding if slide.shapes.title.text == "" else Inches(1.5)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.5), table_top, Inches(9), Inches(6.5)
    )
    table = table_shape.table
    table.rows[0].height = Inches(0.5)
    return table


quarters = {"Q1": [], "Q2": [], "Q3": [], "Q4": []}

for issue in issues:
    summary = issue.fields.summary
    jira_id = issue.key
    link = jira._options["server"] + "/browse/" + issue.key
    labels = issue.fields.labels
    components = []
    for component in issue.fields.components:
        components.append(component.name)
    if not components:
        components = ["No Component"]

    if "PI23/1" in labels:
        quarters["Q1"].append(
            {
                "Summary": summary,
                "JIRA ID": jira_id,
                "Link": link,
                "Label": labels,
                "Components": components,
            }
        )
    if "PI23/2" in labels:
        quarters["Q2"].append(
            {
                "Summary": summary,
                "JIRA ID": jira_id,
                "Link": link,
                "Label": labels,
                "Components": components,
            }
        )
    if "PI23/3" in labels:
        quarters["Q3"].append(
            {
                "Summary": summary,
                "JIRA ID": jira_id,
                "Link": link,
                "Label": labels,
                "Components": components,
            }
        )
    if "PI23/4" in labels:
        quarters["Q4"].append(
            {
                "Summary": summary,
                "JIRA ID": jira_id,
                "Link": link,
                "Label": labels,
                "Components": components,
            }
        )

unique_components = []
for quarter, issues in quarters.items():
    for issue in issues:
        for component in issue["Components"]:
            if component not in unique_components:
                unique_components.append(component)

# Create presentation
prs = pptx.Presentation("./Template.pptx")
padding = Inches(0.2)  # Set padding value

for component in unique_components:
    # Create a new slide and table for each component
    slide = create_slide(prs)
    table = create_table(slide, 2, 5)

    # Set the table width to the width of the slide
    table_width = prs.slide_width - (2 * padding)
    for col in range(4):
        table.columns[col].width = int(table_width / 5)

    # Adjust the width of the last column to include padding
    table.columns[4].width = int(table_width / 5) - padding

    # Set the table position with padding
    table.left = padding

    headers = ["Component", "Q1", "Q2", "Q3", "Q4"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.CENTER

        # Decrease the font size of the column names
        paragraph.runs[0].font.size = pptx.util.Pt(10)

    # Make the component names bold
    cell = table.cell(1, 0)
    cell.text = component
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = pptx.util.Pt(10)
    cell.text_frame.paragraphs[
        0
    ].font.name = "Calibri"  # Set the font type for component

    for col, quarter in enumerate(["Q1", "Q2", "Q3", "Q4"], start=1):
        if start_quarter <= col <= end_quarter:
            summaries = []
            for issue in quarters[quarter]:
                if component in issue["Components"]:
                    summaries.append(
                        (issue["Summary"], issue["JIRA ID"], issue["Link"])
                    )

            cell = table.cell(1, col)
            text_frame = cell.text_frame
            text_frame.clear()  # Clear the existing paragraphs in the cell

            for idx, (summary, jira_id, link) in enumerate(summaries, start=1):
                p = text_frame.add_paragraph()
                p.text = f"• {summary} "
                p.space_after = pptx.util.Pt(6)
                p.level = 0

                # Decrease the font size of the summary
                summary_run = p.runs[0]
                summary_run.font.size = pptx.util.Pt(
                    12
                )  # Change the number to your desired font size
                summary_run.font.name = "Calibri"  # Set the font type for summaries

                r = p.add_run()
                r.text = f"({jira_id})"
                r.hyperlink.address = link

                # Decrease the font size of the issue link
                r.font.size = pptx.util.Pt(
                    10
                )  # Change the number to your desired font size
            else:
                pass


# Save the presentation
prs.save("UseCase2.pptx")
