import configparser
from jira import JIRA
from itertools import product
from pyexpat import features
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.util import Inches
import sys
import os

# Initialize the ConfigParser object
config = configparser.ConfigParser()
config.read("config.ini")

# Get the JIRA credentials from the config file
jira_server = config.get("jira", "server")
jira_username = config.get("jira", "username")
jira_password = config.get("jira", "password")

print("Starting script...")

# Set the JIRA API endpoint URL
jira = JIRA(server=jira_server, basic_auth=(jira_username, jira_password))
# Set the query parameters for the JQL query
jql = 'project = ARTA2C and "Target Release" = "PI23/2"'

# Set the fields to be included in the response
fields = "key,summary,customfield_10241,parent,issuelinks,customfield_10253,customfield_22341"

# Make the API call and fetch the data
issues = jira.search_issues(jql, maxResults=1000, fields=fields)

features = []

print("Processing issues...")
# Loop through the issues and extract the necessary data
for issue in issues:
    # add product
    feature_product = ""
    if issue.fields.customfield_10241:
        for product in issue.fields.customfield_10241:
            feature_product += product.value + ","

    # find parent link
    pe_summary = ""
    pe_benefit = ""
    pe_component = ""
    portfolio_key = ""
    portfolio_epic_url = ""

    # loop through each issue link and get the solution and portfolio epics
    for issuelink in issue.fields.issuelinks:
        if (
            hasattr(issuelink, "inwardIssue")
            and issuelink.inwardIssue.fields.issuetype.name == "Solution Epic"
        ):
            solution_issue = issuelink.inwardIssue
            solution_key = solution_issue.key

            # use JIRA API to get solution issue details
            solution_issue = jira.issue(solution_key)
            solution_issue_links = solution_issue.fields.issuelinks

            for link in solution_issue_links:
                if (
                    hasattr(link, "inwardIssue")
                    and link.inwardIssue.fields.issuetype.name == "Portfolio Epic"
                ):
                    portfolio_issue = link.inwardIssue
                    portfolio_key = portfolio_issue.key

                    # use JIRA API to get portfolio issue details
                    portfolio_issue = jira.issue(portfolio_key)
                    pe_summary = portfolio_issue.fields.summary
                    pe_component = portfolio_issue.fields.components[0].name
                    pe_description = portfolio_issue.fields.description

                    try:
                        pe_benefit = ""
                        start_index = (
                            pe_description.index("Benefit") + len("Benefit") + 1
                        )
                        end_index = pe_description.index("*Description*")
                        pe_benefit = (
                            pe_description[start_index:end_index]
                            .strip()
                            .replace("_x000D_", "")
                        )

                    except:
                        try:
                            if pe_benefit == "":
                                start_index = (
                                    pe_description.index("WE BELIEVE")
                                    + len("WE BELIEVE")
                                    + 1
                                )
                                end_index = pe_description.index("*WILL RESULT IN*")
                                pe_benefit = (
                                    pe_description[start_index:end_index]
                                    .strip()
                                    .replace("_x000D_", "")
                                )
                        except:
                            try:
                                pe_benefit = (
                                    portfolio_issue.fields.customfield_22341.strip()
                                    .replace("_x000D_", "")
                                    .replace("\r\n", " ")
                                )
                            except:
                                pe_benefit = "No benifit found"

    # get feature product
    if hasattr(issue.fields, "customfield_13500"):
        feature_product = issue.fields.customfield_13500.value

    feature_PDD = "No PDD available"
    if hasattr(issue.fields, "customfield_10253"):
        if issue.fields.customfield_10253 != None:
            feature_PDD = issue.fields.customfield_10253

    # append all data
    issue_details = {
        "key": issue.key,
        "summary": issue.fields.summary,
        "feature_url": issue.self,
        "feature_PDD": feature_PDD,
        "product": feature_product,
        "portfolioSummary": pe_summary,
        "benefit": pe_benefit,
        "component": pe_component,
        "portfolio_key": portfolio_key,
        "portfolio_epic_url": portfolio_epic_url,
    }

    features.append(issue_details)

sorted_features = sorted(features, key=lambda x: x["portfolio_key"])

print("Generating Powerpoint file")

if getattr(sys, "frozen", False):
    prs = Presentation(os.path.join(sys._MEIPASS, "./Template.pptx"))
else:
    prs = Presentation("./Template.pptx")

lyt = prs.slide_layouts[1]  # choosing a slide layout
slide1 = prs.slides.add_slide(lyt)  # adding a slide
title = slide1.shapes.title  # assigning a title

# title
title.text = "PI2023.2/RD5 Executive Summary"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(240, 171, 0)
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.italic = True

# Group the features by portfolio_key and product
grouped_features = {}
for feature in sorted_features:
    portfolio_key = feature["portfolio_key"]
    product = feature["product"]
    if portfolio_key not in grouped_features:
        grouped_features[portfolio_key] = {}
    if product not in grouped_features[portfolio_key]:
        grouped_features[portfolio_key][product] = []
    grouped_features[portfolio_key][product].append(feature)

# Calculate the number of rows in the table
num_rows = len(features) + 1  # Add 1 to include the header row

rows_per_slide = 8

# Loop through the grouped features and create slides and tables
for portfolio_key, portfolio_data in grouped_features.items():
    # Loop through the products in the portfolio
    num_rows = len(portfolio_data) + 1
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add table to slide
    table = slide.shapes.add_table(
        2,
        6,
        left=Inches(0.5),
        top=Inches(0.5),
        width=(prs.slide_width - Inches(1)),
        height=(prs.slide_height - Inches(2)),
    ).table

    # Split table into multiple slides if needed
    for i in range(0, len(portfolio_data), rows_per_slide - 1):
        start_row = i
        end_row = min(i + rows_per_slide - 1, len(portfolio_data))

        table.columns[0].width = Inches(2)  # First column
        table.columns[1].width = Inches(1.8)  # Second column
        table.columns[2].width = Inches(3.0)  # Third column
        table.columns[3].width = Inches(1.5)  # Fourth column
        table.columns[4].width = Inches(1.8)  # Fifth column
        table.columns[5].width = Inches(2.6)  # Sixth column

        # Add headers to table
        header_cells = table.rows[0].cells
        header_cells[0].text = "Biz. Capability"
        header_cells[0].text_frame.add_paragraph().text = "Portfolio Epic Component"

        header_cells[1].text = "Summary"
        header_cells[1].text_frame.add_paragraph().text = "Portfolio Epic (SP-) "

        header_cells[2].text = "Benefit"
        header_cells[
            2
        ].text_frame.add_paragraph().text = (
            "Extract from Portfolio Epic Description (first X words)"
        )

        header_cells[3].text = "Go Live Date"
        header_cells[3].text_frame.add_paragraph().text = "Feature PDD"

        header_cells[4].text = "Portfolio Product"
        header_cells[4].text_frame.add_paragraph().text = "Feature Product"

        header_cells[5].text = "Delivered Features"
        header_cells[
            5
        ].text_frame.add_paragraph().text = "Feature Summary (ID incl. link)"

        header_row = table.rows[0]
        header_row.height = Inches(1)

        # Add features to table
        row_count = 1
        portfolio_cells = table.rows[1].cells

        for product in portfolio_data:
            product_cells = table.rows[row_count].cells
            product_cells[4].text_frame.add_paragraph().text = portfolio_data[product][
                0
            ]["product"]

            for feature in portfolio_data[product]:
                feature_cells = table.rows[row_count].cells
                feature_cells[3].text_frame.add_paragraph().text = feature[
                    "feature_PDD"
                ]

                p1 = feature_cells[5].text_frame.add_paragraph()
                p1.text = feature["summary"] + " "
                run = p1.add_run()
                run.text = " " + feature["key"]
                run.font.size = Pt(10)
                link = jira._options["server"] + "/browse/" + feature["key"]
                hyperlink_address = link  # replace with the actual hyperlink address
                run.hyperlink.address = hyperlink_address
                run.hyperlink.anchor = ""
                run.font.underline = True

        portfolio_cells[2].text = portfolio_data[product][0]["benefit"]
        portfolio_cells[0].text = portfolio_data[product][0]["component"]

        p = portfolio_cells[1].text_frame.paragraphs[0]
        p.text = feature["portfolioSummary"] + " "
        run = p.add_run()
        if "_group2" in portfolio_key:
            run.text = portfolio_key.split("_group2")[0]
        else:
            run.text = " " + portfolio_key
        run.font.size = Pt(10)
        link = jira._options["server"] + "/browse/" + portfolio_key
        hyperlink_address = link  # replace with the actual hyperlink address
        run.hyperlink.address = hyperlink_address
        run.hyperlink.anchor = ""
        run.font.underline = True

        # Apply consistent font size to all cells
        font_size = Pt(10)
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = font_size

        for cell in header_cells:
            cell.text_frame.paragraphs[0].font.size = Pt(15)

# print(os.path.dirname(sys.argv[0]))

# prs.save(os.path.dirname(sys.argv[0]) + "/PI2023_2_RD5.pptx")  # saving file
prs.save("UseCase1.pptx")
